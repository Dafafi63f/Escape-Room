#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Genera Presentacion_TFG.pptx a partir de la memoria y figuras del TFG."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

_DOCS = Path(__file__).resolve().parent
_FIGURAS = _DOCS / "Figuras"
_SALIDA = _DOCS / "Entrega" / "Presentacion_TFG.pptx"

_ANCHO = 10.0
_ALTO = 5.625
_MARGEN = 0.42
_TOP_TITULO = 0.38
_ALTO_TITULO = 0.58
_TOP_CUERPO = 1.08
_FONDO_CUERPO = 5.15
_FOOTER_Y = 5.22

_ANCHO_TEXTO_COL = 3.45
_GAP_TEXTO_IMAGEN = 0.12
_IZQ_IMAGEN = _MARGEN + _ANCHO_TEXTO_COL + _GAP_TEXTO_IMAGEN
_ANCHO_IMAGEN_COL = _ANCHO - _MARGEN - _IZQ_IMAGEN
_FRAC_IMAGEN_GRANDE = 0.62

_C_TITULO = RGBColor(0x1A, 0x3A, 0x5C)
_C_ACENTO = RGBColor(0x2E, 0x6B, 0xA4)
_C_ACENTO_CLARO = RGBColor(0xE8, 0xF0, 0xF8)
_C_TEXTO = RGBColor(0x2A, 0x2A, 0x2A)
_C_PIE = RGBColor(0x88, 0x88, 0x88)
_C_OK = RGBColor(0x1B, 0x7A, 0x4A)
_C_PARCIAL = RGBColor(0xB8, 0x6E, 0x00)


def _slides_data() -> list[dict]:
    return [
        {
            "tipo": "portada",
            "titulo": (
                "Diseño y desarrollo de un juego interactivo educativo\n"
                "basado en contenidos del grado MatCAD"
            ),
            "subtitulo": (
                "Daniel Fageda Figueredo · NIU 1601846\n"
                "Tutor: Víctor Navas Portella\n"
                "Universitat Autònoma de Barcelona · v1.0.0 · 2026"
            ),
            "notas": (
                "Buenos días. Presento el TFG: cuestionario gamificado para el grado "
                "MatCAD — 480 preguntas, cinco modos en pygame y herramientas de validación."
            ),
        },
        {
            "titulo": "Índice de la exposición",
            "bullets": [
                "1. Contexto, problema e inspiración (Inka Games → MatCAD)",
                "2. Objetivos, hipótesis, metodología y alcance",
                "3. Banco de preguntas y arquitectura del sistema",
                "4. La aplicación: cinco modos de juego",
                "5. Validación: tests, Monte Carlo y sistema pity",
                "6. Conclusiones y limitaciones",
            ],
            "notas": "15 minutos aprox. Dejar margen para preguntas del tribunal.",
        },
        {
            "titulo": "Contexto: grado MatCAD",
            "bullets": [
                "Matemáticas, programación y ciencia de datos en un plan de 4 cursos.",
                "40 materias · 10 grupos temáticos · evaluación continua.",
                "Necesidad de autoevaluación alineada al plan de estudios.",
                "Gamificación y serious games como complemento formativo.",
            ],
            "notas": "Marco del grado UAB. El proyecto especializa la práctica por materia y etapa.",
        },
        {
            "titulo": "Problema y motivación",
            "bullets": [
                "Practicar por curso, semestre o materia sin montar listados a mano.",
                "Herramientas genéricas sin metadatos curriculares del grado.",
                "Retroalimentación inmediata y mecánicas que incentiven la repetición.",
                "Meta: banco auditable + juego extensible + trazabilidad pedagógica.",
            ],
            "notas": "Complemento a recursos institucionales, no sustituto de AulaWeb o Moodle.",
        },
        {
            "titulo": "Inspiración: Inka Games",
            "layout": "imagen_grande",
            "bullets": [],
            "imagen": "inkagames_gameplay_referencia.png",
            "pie_imagen": "Inka Games — referencia narrativa (uso académico ilustrativo)",
            "notas": (
                "Figura 1 de la memoria. Escape rooms point-and-click: salas, inventario, puertas. "
                "No reproducimos la estética comercial."
            ),
        },
        {
            "titulo": "Inka Games vs. entregable educativo",
            "layout": "dos_columnas",
            "titulo_izq": "Inka Games (referencia)",
            "col_izq": [
                "Entretenimiento comercial",
                "Puzles lógicos y objetos",
                "Narrativa gráfica completa",
                "Sin vínculo curricular",
            ],
            "titulo_der": "TFG MATCAD (entregable)",
            "col_der": [
                "Autoevaluación formativa",
                "Preguntas A–D del banco",
                "Mecánicas jugables; guion visual pendiente",
                "40 materias con metadatos",
            ],
            "notas": "Tabla 1 de la memoria. Priorizamos el núcleo evaluable antes de mostrar capturas.",
        },
        {
            "titulo": "Adaptación: escape room MatCAD",
            "layout": "imagenes_lado",
            "bullets": [],
            "imagenes": ["tfg_escape_sala_puertas.png", "tfg_escape_pregunta.png"],
            "pies_imagen": ["Sala 1 — elige puerta", "Pregunta del banco con inventario"],
            "notas": "Izq.: sala con 3 puertas · Der.: pregunta A–D con inventario. Capturas pygame (semilla fija).",
        },
        {
            "titulo": "Alcance y entregable",
            "layout": "metricas",
            "metricas": [
                ("5", "modos de juego"),
                ("480", "preguntas revisadas"),
                ("578", "tests automáticos"),
                ("2", "paquetes zip"),
            ],
            "bullets": [
                "Cuestionario gamificado + banco estructurado + scripts de mantenimiento.",
                "Distribución portable (completo) y mínima (solo CSV).",
            ],
            "notas": "Entregable v1.0.0 cerrado en junio 2026.",
        },
        {
            "titulo": "Objetivos del proyecto",
            "layout": "objetivos",
            "filas": [
                ("General", "Juego educativo con retos del grado MatCAD", "Cumplido"),
                ("OE1", "Narrativa interactiva (escape room)", "Parcial"),
                ("OE2", "Retos por materias (480 ítems)", "Cumplido"),
                ("OE3", "Validación automática A–D", "Cumplido"),
                ("OE4", "Interfaz gráfica pygame-ce", "Cumplido"),
                ("OE5", "Valor formativo contrastado", "Parcial"),
            ],
            "notas": "OE1 parcial en capa narrativa visual; OE5 sin piloto con usuarios.",
        },
        {
            "titulo": "Hipótesis de trabajo",
            "bullets": [
                "H1: banco segmentado → autoevaluación por filtros (modo libre).",
                "H2: pipeline automatizado → inconsistencias reproducibles.",
                "H3: histórico 8 818 notas → ponderación en modo historia.",
                "Monte Carlo: el azar no aprueba (nota media 2,5/10).",
            ],
            "notas": "H1–H3 contrastadas con datos del sistema. Motivación → piloto futuro.",
        },
        {
            "titulo": "Metodología en cuatro fases",
            "bullets": [
                "Fase 1 — Análisis: plan MatCAD, esquema del banco, requisitos.",
                "Fase 2 — Implementación: Python, pygame-ce, 5 modos, scripts.",
                "Fase 3 — Pruebas: revisión 480/480, 578 tests, CI GitHub Actions.",
                "Fase 4 — Evaluación: Monte Carlo, pity, memoria y limitaciones.",
            ],
            "notas": "Desarrollo incremental con control de versiones Git.",
        },
        {
            "titulo": "Banco de preguntas",
            "bullets": [
                "480 ítems · 40 materias · 12/materia (2FT…2DC).",
                "160 Fácil · 160 Media · 160 Difícil.",
                "240 Teoría · 240 Cálculo.",
                "Revisión manual 480/480 · 0 duplicados.",
                "Metadatos: curso, semestre, grupo G1–G10, nivel, temática.",
            ],
            "notas": "Banco cerrado junio 2026. Ampliado (960) opcional y etiquetado aparte.",
        },
        {
            "titulo": "Arquitectura en cinco capas",
            "bullets": [
                "1. Lanzador — juego_grafico.py (pygame-ce)",
                "2. Modos — libre, historia, resistencia, escape, feedback",
                "3. Motor — Comun/ (reglas, vidas, puntuación, informes)",
                "4. Datos — CSV/JSON en Data/Banco/ y Data/Juego/",
                "5. Mantenimiento — Files/mantenimiento.py (fuera del juego)",
            ],
            "notas": "Modularidad: nuevos modos sin romper el motor ni el banco.",
        },
        {
            "titulo": "Cinco modos operativos",
            "layout": "columna",
            "frac_imagen": _FRAC_IMAGEN_GRANDE,
            "bullets": [
                "Libre — filtros multidimensionales.",
                "Historia — examen según histórico.",
                "Resistencia — rachas, apuestas, eventos.",
                "Escape — salas, tienda, inventario.",
                "Feedback — mejora continua del banco.",
            ],
            "imagen": "tfg_menu_principal.png",
            "notas": "Barra superior: examen del día, estadísticas, opciones.",
        },
        {
            "titulo": "Modo historia (validación H3)",
            "layout": "columna",
            "frac_imagen": _FRAC_IMAGEN_GRANDE,
            "bullets": [
                "8 818 registros de calificaciones MatCAD.",
                "Índice de dificultad por materia.",
                "5 presets: repaso, simulacro, examen asignatura…",
                "Examen dirigido tras fallos en la sesión.",
                "Examen del día: 24 preguntas compartidas (semilla de fecha).",
            ],
            "imagen": "tfg_historia_carrusel.png",
            "notas": (
                "Carrusel de presets con prioridad histórica. "
                "Ejemplo materias exigentes: Càlcul DV, Probabilitat, Anàlisi Complexa."
            ),
        },
        {
            "titulo": "Gamificación: resistencia y escape",
            "layout": "imagenes_lado",
            "bullets": [
                "Resistencia: escalada, maldiciones, power-ups.",
                "Escape: 5–50 salas, economía, jefe final.",
                "Inventario compartido entre modos arcade.",
            ],
            "imagenes": ["tfg_resistencia_partida.png", "tfg_escape_tienda.png"],
            "pies_imagen": [
                "Modo resistencia — partida en curso",
                "Escape room — tienda entre salas",
            ],
            "notas": "Mecánicas inspiradas en Inka; contenido siempre del banco MatCAD.",
        },
        {
            "titulo": "Validación y control de calidad",
            "layout": "cuatro_bloques",
            "bloques": [
                ("578", "tests + CI"),
                ("0", "duplicados"),
                ("480", "ítems revisados"),
                ("50 000", "réplicas Monte Carlo"),
            ],
            "bullets": [
                "mantenimiento.py validar · auditar-distractores · duplicados.py",
                "github.com/Dafafi63f/Escape-Room",
            ],
            "notas": "Integridad verificada en cada push (GitHub Actions).",
        },
        {
            "titulo": "Simulación Monte Carlo",
            "layout": "imagenes_lado",
            "bullets": [
                "Respuestas al azar (p = 1/4 por ítem).",
                "Nota media ≈ 2,5/10 en examen de 20 preguntas.",
                "Fracción de aciertos ≈ 25 % (50 000 réplicas).",
                "Aprobar por azar: estadísticamente inviable.",
            ],
            "imagenes": [
                "monte_carlo_histograma_notas.png",
                "monte_carlo_convergencia.png",
            ],
            "pies_imagen": [
                "Histograma de notas",
                "Convergencia de la media",
            ],
            "notas": "Valida el motor de corrección; detalle en memoria §5.7.",
        },
        {
            "titulo": "Sistema pity (equidad en partidas largas)",
            "layout": "dos_columnas",
            "titulo_izq": "Sin pity",
            "col_izq": [
                "15,6 % partidas sin descanso",
                "Racha p95: 29 salas",
                "Mayor frustración aleatoria",
            ],
            "titulo_der": "Con pity (implementado)",
            "col_der": [
                "0 % partidas sin descanso",
                "Racha p95: 4 salas",
                "Motivación sostenida",
            ],
            "imagenes": [
                "pity_comparacion_descanso.png",
                "pity_distribucion_primer_descanso.png",
            ],
            "layout_extra": "imagen_pie",
            "frac_imagen_pie": 0.62,
            "notas": "10 000 réplicas · 30 salas · simulacion_pity.py (memoria §5.8).",
        },
        {
            "titulo": "Contribución y conclusiones",
            "bullets": [
                "De listado genérico a herramienta con criterio didáctico explícito.",
                "Banco trazable + cinco modos jugables + arquitectura extensible.",
                "Paquetes portables listos para distribución (Python + pygame-ce).",
                "Base sólida para piloto, narrativa gráfica e integración Moodle.",
            ],
            "notas": "Competencias: Python, diseño de datos, ingeniería incremental.",
        },
        {
            "titulo": "Limitaciones y trabajo futuro",
            "layout": "dos_columnas",
            "titulo_izq": "Limitaciones",
            "col_izq": [
                "Sin piloto con usuarios",
                "Narrativa gráfica incompleta",
                "Una materia por pregunta",
                "Sin prerrequisitos en CSV",
            ],
            "titulo_der": "Trabajo futuro",
            "col_der": [
                "Piloto 15–20 estudiantes (SUS)",
                "Materias_relacionadas / Prerequisitos",
                "Escape room narrativo completo",
                "Moodle · panel de feedback",
            ],
            "notas": "Diseño de piloto ya descrito en memoria §7.",
        },
        {
            "tipo": "cierre",
            "titulo": "Gracias — ¿Preguntas?",
            "subtitulo": (
                "Daniel Fageda Figueredo · NIU 1601846\n"
                "Tutor: Víctor Navas Portella · UAB\n"
                "github.com/Dafafi63f/Escape-Room"
            ),
            "notas": "Agradecer al tutor y al tribunal.",
        },
    ]


# --- Estilos y geometría -----------------------------------------------------

def _tamano_fuente(n: int) -> int:
    if n <= 4:
        return 14
    if n <= 5:
        return 13
    return 12


def _estilo_titulo(p, *, portada: bool = False, size: int | None = None) -> None:
    p.font.bold = True
    p.font.color.rgb = _C_TITULO
    p.font.size = Pt(size or (26 if portada else 22))
    p.alignment = PP_ALIGN.CENTER if portada else PP_ALIGN.LEFT


def _estilo_cuerpo(p, *, portada: bool = False, size: int = 14) -> None:
    p.font.color.rgb = _C_TEXTO
    p.font.size = Pt(16 if portada else size)
    p.alignment = PP_ALIGN.CENTER if portada else PP_ALIGN.LEFT


def _marco(tf, *, ajustar: bool = True) -> None:
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if ajustar:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _bullets(tf, items: list[str], *, size: int | None = None) -> None:
    tf.clear()
    pt = size or _tamano_fuente(len(items))
    for i, t in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.level = 0
        p.space_after = Pt(3)
        p.line_spacing = 1.05
        for r in p.runs:
            r.font.size = Pt(pt)
            r.font.color.rgb = _C_TEXTO


def _dim_img(nombre: str, max_w: float, max_h: float) -> tuple[float, float]:
    try:
        from PIL import Image

        with Image.open(_FIGURAS / nombre) as im:
            pw, ph = im.size
    except Exception:
        return max_w, max_h * 0.72
    ar = ph / pw if pw else 1.0
    w, h = max_w, max_w * ar
    if h > max_h:
        h = max_h
        w = h / ar if ar else max_w
    return w, h


def _img(
    slide,
    nombre: str,
    left: float,
    top: float,
    max_w: float,
    max_h: float,
    *,
    centrar: bool = False,
    alinear_arriba: bool = False,
) -> tuple[float, float, float, float]:
    """Inserta imagen ajustada. Devuelve (left, top, width, height) en pulgadas."""
    ruta = _FIGURAS / nombre
    if not ruta.is_file():
        return left, top, max_w, max_h
    w, h = _dim_img(nombre, max_w, max_h)
    if alinear_arriba:
        top_adj = top
    else:
        top_adj = top + max(0.0, (max_h - h) / 2)
    left_adj = (_ANCHO - w) / 2 if centrar else left
    slide.shapes.add_picture(
        str(ruta),
        Inches(left_adj),
        Inches(top_adj),
        width=Inches(w),
        height=Inches(h),
    )
    return left_adj, top_adj, w, h


def _decoracion(slide, n: int, total: int) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(_ANCHO), Inches(0.14)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _C_ACENTO
    bar.line.fill.background()

    pie = slide.shapes.add_textbox(
        Inches(_MARGEN), Inches(_FOOTER_Y), Inches(_ANCHO - 2 * _MARGEN), Inches(0.28)
    )
    p = pie.text_frame.paragraphs[0]
    p.text = f"TFG MatCAD — Daniel Fageda Figueredo          {n}/{total}"
    p.font.size = Pt(9)
    p.font.color.rgb = _C_PIE
    p.alignment = PP_ALIGN.RIGHT


def _titulo_slide(slide, texto: str) -> None:
    box = slide.shapes.add_textbox(
        Inches(_MARGEN),
        Inches(_TOP_TITULO),
        Inches(_ANCHO - 2 * _MARGEN),
        Inches(_ALTO_TITULO),
    )
    _marco(box.text_frame, ajustar=False)
    p = box.text_frame.paragraphs[0]
    p.text = texto
    _estilo_titulo(p)


# --- Tipos de diapositiva ----------------------------------------------------

def _slide_portada(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(_ANCHO), Inches(0.18)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _C_TITULO
    bar.line.fill.background()

    fondo = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(4.6),
        Inches(_ANCHO),
        Inches(1.025),
    )
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = _C_ACENTO_CLARO
    fondo.line.fill.background()

    b1 = slide.shapes.add_textbox(Inches(0.55), Inches(1.35), Inches(8.9), Inches(2.0))
    _marco(b1.text_frame, ajustar=False)
    p = b1.text_frame.paragraphs[0]
    p.text = data["titulo"]
    _estilo_titulo(p, portada=True)

    b2 = slide.shapes.add_textbox(Inches(0.55), Inches(3.55), Inches(8.9), Inches(1.1))
    p2 = b2.text_frame.paragraphs[0]
    p2.text = data["subtitulo"]
    _estilo_cuerpo(p2, portada=True)

    if data.get("notas"):
        slide.notes_slide.notes_text_frame.text = data["notas"]


def _slide_solo_texto(slide, data: dict) -> None:
    h = _FONDO_CUERPO - _TOP_CUERPO
    box = slide.shapes.add_textbox(
        Inches(_MARGEN), Inches(_TOP_CUERPO), Inches(_ANCHO - 2 * _MARGEN), Inches(h)
    )
    _marco(box.text_frame)
    _bullets(box.text_frame, data.get("bullets", []))


def _columna_ancho_texto_imagen(data: dict) -> tuple[float, float, float]:
    """Devuelve (ancho_texto, x_imagen, ancho_imagen) en pulgadas."""
    ancho_util = _ANCHO - 2 * _MARGEN
    frac = data.get("frac_imagen")
    if frac is not None:
        w_img = ancho_util * float(frac)
        w_txt = ancho_util - w_img - _GAP_TEXTO_IMAGEN
        return w_txt, _MARGEN + w_txt + _GAP_TEXTO_IMAGEN, w_img
    return _ANCHO_TEXTO_COL, _IZQ_IMAGEN, _ANCHO_IMAGEN_COL


def _slide_columna(slide, data: dict) -> None:
    h = _FONDO_CUERPO - _TOP_CUERPO
    w_txt, x_img, w_img = _columna_ancho_texto_imagen(data)
    box = slide.shapes.add_textbox(
        Inches(_MARGEN), Inches(_TOP_CUERPO), Inches(w_txt), Inches(h)
    )
    _marco(box.text_frame)
    _bullets(box.text_frame, data.get("bullets", []), size=13)
    if img := data.get("imagen"):
        _img(slide, img, x_img, _TOP_CUERPO, w_img, h)


def _slide_imagenes_lado(slide, data: dict) -> None:
    """Dos capturas lado a lado (p. ej. puertas | pregunta escape)."""
    bullets = data.get("bullets", [])
    h_txt = 0.0
    if bullets:
        h_txt = 0.72 if len(bullets) <= 1 else 1.05
        box = slide.shapes.add_textbox(
            Inches(_MARGEN),
            Inches(_TOP_CUERPO),
            Inches(_ANCHO - 2 * _MARGEN),
            Inches(h_txt),
        )
        _marco(box.text_frame)
        _bullets(box.text_frame, bullets, size=13)

    pies = data.get("pies_imagen") or []
    pie_h = 0.28 if pies else 0.0
    top = _TOP_CUERPO + (h_txt + 0.04 if h_txt else 0.0)
    h_img = _FONDO_CUERPO - top - pie_h - 0.06

    nombres = data.get("imagenes", [])
    if not nombres:
        return
    gap = 0.14
    w_each = (_ANCHO - 2 * _MARGEN - gap * (len(nombres) - 1)) / len(nombres)
    for i, nombre in enumerate(nombres):
        x = _MARGEN + i * (w_each + gap)
        _, top_img, _, img_h = _img(
            slide, nombre, x, top, w_each, h_img, alinear_arriba=True
        )
        if i < len(pies):
            cap = slide.shapes.add_textbox(
                Inches(x),
                Inches(top_img + img_h + 0.05),
                Inches(w_each),
                Inches(pie_h),
            )
            p = cap.text_frame.paragraphs[0]
            p.text = pies[i]
            p.font.size = Pt(10)
            p.font.color.rgb = _C_PIE
            p.alignment = PP_ALIGN.CENTER


def _slide_imagen_grande(slide, data: dict) -> None:
    """Una sola imagen centrada y ampliada (p. ej. Inka vs escape en slides distintas)."""
    bullets = data.get("bullets", [])
    h_txt = 0.0
    if bullets:
        h_txt = 0.55 if len(bullets) <= 1 else 0.95
        box = slide.shapes.add_textbox(
            Inches(_MARGEN),
            Inches(_TOP_CUERPO),
            Inches(_ANCHO - 2 * _MARGEN),
            Inches(h_txt),
        )
        _marco(box.text_frame)
        _bullets(box.text_frame, bullets, size=13)

    top = _TOP_CUERPO + (h_txt + 0.05 if h_txt else 0.0)
    pie_h = 0.24 if data.get("pie_imagen") else 0.0
    h_img = _FONDO_CUERPO - top - pie_h - 0.03
    max_w = _ANCHO - 2 * _MARGEN

    if img := data.get("imagen"):
        _, top_img, _, img_h = _img(
            slide, img, _MARGEN, top, max_w, h_img, centrar=True, alinear_arriba=True
        )
        top_pie = top_img + img_h + 0.05
    else:
        top_pie = top

    if pie := data.get("pie_imagen"):
        cap = slide.shapes.add_textbox(
            Inches(_MARGEN),
            Inches(top_pie),
            Inches(_ANCHO - 2 * _MARGEN),
            Inches(pie_h),
        )
        p = cap.text_frame.paragraphs[0]
        p.text = pie
        p.font.size = Pt(9)
        p.font.color.rgb = _C_PIE
        p.alignment = PP_ALIGN.CENTER


def _slide_imagenes_abajo(slide, data: dict) -> None:
    h_txt = 1.45
    box = slide.shapes.add_textbox(
        Inches(_MARGEN),
        Inches(_TOP_CUERPO),
        Inches(_ANCHO - 2 * _MARGEN),
        Inches(h_txt),
    )
    _marco(box.text_frame)
    _bullets(box.text_frame, data.get("bullets", []), size=13)

    top = _TOP_CUERPO + h_txt + 0.1
    h_img = _FONDO_CUERPO - top
    nombres = data.get("imagenes", [])
    if not nombres:
        return
    gap = 0.15
    w_each = (_ANCHO - 2 * _MARGEN - gap * (len(nombres) - 1)) / len(nombres)
    x = _MARGEN
    for nombre in nombres:
        _img(slide, nombre, x, top, w_each, h_img)
        x += w_each + gap


def _slide_dos_columnas(slide, data: dict) -> None:
    h = _FONDO_CUERPO - _TOP_CUERPO
    w_col = (_ANCHO - 2 * _MARGEN - 0.25) / 2
    extra_img = data.get("layout_extra") == "imagen_pie"
    frac_pie = float(data.get("frac_imagen_pie", 0.62))
    h_cols = h * (1.0 - frac_pie) if extra_img else h
    h_cols = max(h_cols, 1.35 if extra_img else h)

    for i, (tit, key) in enumerate(
        (("titulo_izq", "col_izq"), ("titulo_der", "col_der"))
    ):
        x = _MARGEN + i * (w_col + 0.25)
        th = slide.shapes.add_textbox(Inches(x), Inches(_TOP_CUERPO), Inches(w_col), Inches(0.35))
        p = th.text_frame.paragraphs[0]
        p.text = data.get(tit, "")
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = _C_ACENTO

        body = slide.shapes.add_textbox(
            Inches(x), Inches(_TOP_CUERPO + 0.38), Inches(w_col), Inches(h_cols - 0.38)
        )
        _marco(body.text_frame)
        _bullets(body.text_frame, data.get(key, []), size=12)

    if extra_img:
        nombres = data.get("imagenes") or (
            [data["imagen"]] if data.get("imagen") else []
        )
        if nombres:
            top = _TOP_CUERPO + h_cols + 0.06
            h_pie = _FONDO_CUERPO - top
            gap = 0.12
            w_each = (_ANCHO - 2 * _MARGEN - gap * (len(nombres) - 1)) / len(nombres)
            x = _MARGEN
            for nombre in nombres:
                _img(slide, nombre, x, top, w_each, h_pie)
                x += w_each + gap


def _slide_metricas(slide, data: dict) -> None:
    metricas = data.get("metricas", [])
    n = len(metricas)
    gap = 0.12
    w = (_ANCHO - 2 * _MARGEN - gap * (n - 1)) / max(n, 1)
    y = _TOP_CUERPO
    h_box = 1.05
    for i, (valor, etiqueta) in enumerate(metricas):
        x = _MARGEN + i * (w + gap)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h_box),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = _C_ACENTO_CLARO
        rect.line.color.rgb = _C_ACENTO

        tb = slide.shapes.add_textbox(Inches(x), Inches(y + 0.08), Inches(w), Inches(h_box - 0.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = valor
        p1.alignment = PP_ALIGN.CENTER
        p1.font.bold = True
        p1.font.size = Pt(22)
        p1.font.color.rgb = _C_TITULO
        p2 = tf.add_paragraph()
        p2.text = etiqueta
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(10)
        p2.font.color.rgb = _C_TEXTO

    if bullets := data.get("bullets"):
        box = slide.shapes.add_textbox(
            Inches(_MARGEN),
            Inches(y + h_box + 0.2),
            Inches(_ANCHO - 2 * _MARGEN),
            Inches(_FONDO_CUERPO - y - h_box - 0.25),
        )
        _marco(box.text_frame)
        _bullets(box.text_frame, bullets, size=13)


def _slide_objetivos(slide, data: dict) -> None:
    filas = data.get("filas", [])
    y0 = _TOP_CUERPO
    row_h = min(0.52, (_FONDO_CUERPO - y0 - 0.05) / max(len(filas), 1))
    cols = [0.7, 5.5, 1.5]
    xs = [_MARGEN, _MARGEN + cols[0] + 0.08, _MARGEN + cols[0] + cols[1] + 0.16]

    for i, (cod, desc, estado) in enumerate(filas):
        y = y0 + i * row_h
        if i % 2 == 0:
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(_MARGEN),
                Inches(y),
                Inches(_ANCHO - 2 * _MARGEN),
                Inches(row_h - 0.04),
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = _C_ACENTO_CLARO
            bg.line.fill.background()

        for j, (txt, w) in enumerate(((cod, cols[0]), (desc, cols[1]), (estado, cols[2]))):
            tb = slide.shapes.add_textbox(Inches(xs[j]), Inches(y + 0.06), Inches(w), Inches(row_h - 0.08))
            p = tb.text_frame.paragraphs[0]
            p.text = txt
            p.font.size = Pt(11 if j < 2 else 10)
            p.font.bold = j == 0
            if j == 2:
                p.font.color.rgb = _C_OK if "Cumplido" in txt else _C_PARCIAL
            else:
                p.font.color.rgb = _C_TEXTO


def _slide_cuatro_bloques(slide, data: dict) -> None:
    bloques = data.get("bloques", [])
    gap = 0.1
    w = (_ANCHO - 2 * _MARGEN - gap * 3) / 4
    y = _TOP_CUERPO
    h = 0.95
    for i, (valor, etiqueta) in enumerate(bloques):
        x = _MARGEN + i * (w + gap)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = _C_ACENTO
        rect.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(x), Inches(y + 0.1), Inches(w), Inches(h - 0.12))
        tf = tb.text_frame
        p1 = tf.paragraphs[0]
        p1.text = valor
        p1.alignment = PP_ALIGN.CENTER
        p1.font.bold = True
        p1.font.size = Pt(18)
        p1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p2 = tf.add_paragraph()
        p2.text = etiqueta
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(9)
        p2.font.color.rgb = RGBColor(0xE8, 0xF0, 0xF8)

    if bullets := data.get("bullets"):
        box = slide.shapes.add_textbox(
            Inches(_MARGEN),
            Inches(y + h + 0.22),
            Inches(_ANCHO - 2 * _MARGEN),
            Inches(_FONDO_CUERPO - y - h - 0.3),
        )
        _marco(box.text_frame)
        _bullets(box.text_frame, bullets, size=13)


def _slide_cierre(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(_ANCHO), Inches(_ALTO)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _C_TITULO
    bar.line.fill.background()

    b1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9.0), Inches(1.0))
    p = b1.text_frame.paragraphs[0]
    p.text = data["titulo"]
    _estilo_titulo(p, portada=True, size=30)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    b2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9.0), Inches(1.4))
    p2 = b2.text_frame.paragraphs[0]
    p2.text = data.get("subtitulo", "")
    _estilo_cuerpo(p2, portada=True)
    p2.font.color.rgb = RGBColor(0xD0, 0xE4, 0xF5)

    if data.get("notas"):
        slide.notes_slide.notes_text_frame.text = data["notas"]


def _slide_contenido(prs: Presentation, data: dict, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _decoracion(slide, n, total)
    _titulo_slide(slide, data["titulo"])

    layout = data.get("layout", "solo_texto")
    {
        "columna": _slide_columna,
        "imagen_grande": _slide_imagen_grande,
        "imagenes_lado": _slide_imagenes_lado,
        "imagenes_abajo": _slide_imagenes_abajo,
        "dos_columnas": _slide_dos_columnas,
        "metricas": _slide_metricas,
        "objetivos": _slide_objetivos,
        "cuatro_bloques": _slide_cuatro_bloques,
    }.get(layout, _slide_solo_texto)(slide, data)

    if data.get("notas"):
        slide.notes_slide.notes_text_frame.text = data["notas"]


def generar_presentacion(destino: Path | None = None) -> Path:
    destino = destino or _SALIDA
    destino.parent.mkdir(parents=True, exist_ok=True)
    slides = _slides_data()
    total = len(slides)

    prs = Presentation()
    prs.slide_width = Inches(_ANCHO)
    prs.slide_height = Inches(_ALTO)

    contenido_idx = 0
    for data in slides:
        tipo = data.get("tipo", "contenido")
        if tipo == "portada":
            _slide_portada(prs, data)
        elif tipo == "cierre":
            _slide_cierre(prs, data)
        else:
            contenido_idx += 1
            _slide_contenido(prs, data, contenido_idx, total - 2)

    prs.save(str(destino))
    return destino


def main() -> None:
    ruta = generar_presentacion()
    print("OK:", str(ruta).encode("ascii", "replace").decode(), f"({len(_slides_data())} diapositivas)")


if __name__ == "__main__":
    main()
